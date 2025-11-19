import io
import pdfplumber
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract

class ParserService:
    async def parse_pdf(self, content: bytes) -> str:
        """Извлечь текст из PDF"""
        text = ""
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    
    async def parse_pptx(self, content: bytes) -> str:
        """Извлечь текст из PPTX"""
        text = ""
        prs = Presentation(io.BytesIO(content))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    
    async def parse_docx(self, content: bytes) -> str:
        """Извлечь текст из DOCX"""
        doc = Document(io.BytesIO(content))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    
    async def parse_image(self, content: bytes) -> str:
        """Извлечь текст из изображения (OCR)"""
        image = Image.open(io.BytesIO(content))
        text = pytesseract.image_to_string(image, lang='rus+eng')
        return text.strip()
